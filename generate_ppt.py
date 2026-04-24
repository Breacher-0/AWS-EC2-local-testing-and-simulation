from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # MAX READABILITY COLORS (High Contrast)
    COLOR_BG = RGBColor(255, 255, 255)    # Pure White
    COLOR_TEXT = RGBColor(0, 0, 0)        # Pure Black
    COLOR_ACCENT = RGBColor(0, 51, 102)   # Deep Navy Blue (For titles)
    COLOR_TABLE_HEADER = RGBColor(240, 240, 240) # Light Gray

    def apply_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_title_slide():
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        apply_background(slide)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Cloud Simulation & Software Validation Demo"
        subtitle.text = "Ahmed Firas Mahmoud Khalil\nStudent ID: U23LYAZ801\n\nSubject: Building and Testing a Localized AWS Ecosystem"
        
        # Style
        title.text_frame.paragraphs[0].font.color.rgb = COLOR_ACCENT
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(40)
        
        subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT
        subtitle.text_frame.paragraphs[0].font.size = Pt(28)

    def add_text_slide(title_text, lines):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        apply_background(slide)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLOR_ACCENT
        title.text_frame.paragraphs[0].font.bold = True
        
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.word_wrap = True
        
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = line
            p.font.color.rgb = COLOR_TEXT
            p.font.size = Pt(22)
            p.space_after = Pt(12)

    # --- CONTENT GENERATION (Literal copy from user guide) ---

    add_title_slide()

    # 1. Core Concept
    add_text_slide("1. The Core Concept: 'The Magic Box'", [
        "To understand this project, you must first understand the environment.",
        "In professional engineering, we don’t want to test directly on the internet where it costs money and is slow. Instead, we use a 'Magic Box.'",
        "• Docker: An isolated, tiny computer inside your laptop. It carves out a slice of your CPU and memory separate from your personal files."
    ])

    add_text_slide("1. The Core Concept (Continued)", [
        "• LocalStack: A simulator inside the Docker 'Magic Box' that 'pretends' to be the AWS cloud. It looks, acts, and talks exactly like the real AWS.",
        "• The Tunnel (Port 4566): We created a digital tunnel (Port 4566) that connects your computer to the simulator inside the box. Every piece of code we wrote sends its data through this tunnel."
    ])

    # 2. Meet the Characters
    add_text_slide("2. Meet the Characters (The AWS Services)", [
        "We are simulating three specific cloud tools, each with a unique job:",
        "• S3 (Simple Storage Service): A giant, programmable hard drive. We created a 'Bucket' (a folder) called university-docs to store the files you upload.",
        "• Lambda (The Worker): This is 'Serverless' computing. A worker who stays asleep until an alarm goes off. When a file is uploaded, the worker wakes up, validates, and goes back to sleep."
    ])

    add_text_slide("2. Meet the Characters (Continued)", [
        "• CloudWatch (The Logbook): This is a digital notebook.",
        "• Every time the Worker (Lambda) wakes up, they write a note here so we can verify their work later."
    ])

    # 3. The Invisible Conversation
    add_text_slide("3. The Invisible Conversation (The Plumbing)", [
        "1. The Courier (AWS SDK): Our React website uses a special library called the AWS SDK. Think of this as a highly trained Courier. We gave the Courier a specific address: http://localhost:4566.",
        "2. The Handshake: When you click 'Upload,' the website hands the file to the Courier. The Courier packages it into a specific format that cloud systems understand."
    ])

    add_text_slide("3. The Invisible Conversation (Continued)", [
        "3. The Receptionist: Inside the Magic Box, LocalStack acts as a Receptionist.",
        "• It catches the file, identifies it as an 'S3 Storage' request, and places it in the storage room."
    ])

    # 4. The Reaction
    add_text_slide("4. The Reaction: 'The Silent Alarm'", [
        "The most powerful part of this project is the Event-Driven Trigger.",
        "• The Setup: We wrote a blueprint (the serverless.yml file) that installed a sensor in the S3 storage room.",
        "• The Trigger: The moment a file hits the floor of the S3 bucket, a Silent Alarm goes off."
    ])

    add_text_slide("4. The Reaction (Continued)", [
        "• The Lambda Reaction: The Alarm immediately wakes up the Lambda worker.",
        "• The worker receives a 'Note' (a JSON Event Object) telling it the file's name and size.",
        "• The worker then writes its validation message in the CloudWatch Logbook."
    ])

    # 5. The Dashboard
    add_text_slide("5. The Dashboard (The React Frontend)", [
        "• State Management: The app is a 'State Machine.' It tracks whether it is currently 'Idle,' 'Uploading,' or if it has reached a 'Success' state.",
        "• This allows the UI to change (e.g., showing a progress spinner) in real-time.",
        "• Visual Excellence: We used Vanilla CSS to create a premium, dark-mode design. Used 'Clean Class Names' (like .upload-btn)."
    ])

    # 6. The Audit
    add_text_slide("6. The Audit (Playwright Validation)", [
        "As a Lead QA Engineer, we don't just hope it works; we prove it.",
        "1. Automated Interaction: The Robot opens the browser, finds the upload button, and 'injects' a test file.",
        "2. API Interrogation: The Robot talks directly to the LocalStack API and asks: 'Show me the list of files. Is test-file.txt there?'"
    ])

    add_text_slide("6. The Audit (Continued)", [
        "3. Log Verification: The Robot then checks the CloudWatch Logbook to ensure the Lambda worker actually recorded a 'Validation Event.'"
    ])

    # 7. The Table
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_background(slide)
    title = slide.shapes.title
    title.text = "7. File-by-File Purpose Breakdown"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_ACCENT

    rows, cols = 11, 2
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(6)).table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(6.8)

    data = [
        ("File Name", "Purpose"),
        ("docker-compose.yml", "Instructions to build and start the 'Magic Box' (LocalStack)."),
        ("package.json (Root)", "The remote control to start the demo with one command."),
        ("serverless.yml", "Blueprint for S3 buckets and Lambda workers."),
        ("handler.js", "The actual code the Lambda worker runs (The Logic)."),
        ("App.jsx", "Frontend dashboard logic and SDK connection."),
        ("App.css", "Premium styling and user interface layout."),
        ("main.jsx", "The 'Ignition' file that starts the React engine."),
        ("upload.spec.js", "The script for the Playwright Robot Auditor."),
        ("playwright.config.js", "Master settings for the Robot Auditor."),
        ("README.md", "High-level guide for other engineers.")
    ]

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_TEXT
            if r == 0:
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_TABLE_HEADER

    # 8. Summary
    add_text_slide("8. Summary for the Demo", [
        "\"This project demonstrates a fully automated, self-validating cloud pipeline. By simulating AWS locally via Docker and LocalStack, we have created a risk-free environment to test Event-Driven Architecture. We have successfully linked a React dashboard to a Lambda function via S3, and we have used Playwright to provide a 100% automated proof of quality.\""
    ])

    prs.save('Ahmed_Khalil_Software_Validation.pptx')

if __name__ == "__main__":
    create_presentation()
